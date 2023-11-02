

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
logo_path = '/home/godfrey/Documents/376774413_300965799294806_6095592136693544045_n.jpg'  # Replace with the actual path to your logo image

# Create a prs
prs = Presentation()
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
 
# Slide 1: Cover Page
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide

# Add a full-screen image to the cover page
left_inch = Inches(0)
top_inch = Inches(0)
width_inch = Inches(10)
height_inch = Inches(7.5)
image = slide_1.shapes.add_picture(logo_path, left_inch, top_inch, width_inch, height_inch)

# Add the title text
#                                         xxxx,(0=top, 7=down), (0=left, 7=right),xxxxx
title = slide_1.shapes.add_textbox(Inches(0.5), Inches(7), Inches(2), Inches(1))
title_frame = title.text_frame
title_frame.text = "Climate Report" 
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add the name and date at the bottom left
name_date = slide_1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(4), Inches(1))
name_date_frame = name_date.text_frame
name_date_frame.text = "Prepared by: Your Name\nDate: August 14, 2023"
name_date_frame.paragraphs[0].font.size = Pt(18)

# ===================================================================== 
# Slide 2: Body with Two Panels
slide_2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank Slide

# Add panel 1 with a heading
left_inch = Inches(0.5)
top_inch = Inches(1)
width_inch = Inches(4)
height_inch = Inches(3)
panel_1 = slide_2.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
text_frame_1 = panel_1.text_frame

# Add a heading to panel 1
heading_1 = text_frame_1.add_paragraph()
heading_1.text = "Panel 1 Heading"
heading_1.font.size = Inches(0.5)

# Add content to panel 1 (if needed)
content_1 = text_frame_1.add_paragraph()
content_1.text = "This is the content of panel 1."

# Add panel 2 with a heading
left_inch = Inches(5.5)
top_inch = Inches(1)
width_inch = Inches(4)
height_inch = Inches(3)
panel_2 = slide_2.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
text_frame_2 = panel_2.text_frame

# Add a heading to panel 2
heading_2 = text_frame_2.add_paragraph()
heading_2.text = "Panel 2 Heading"
heading_2.font.size = Inches(0.5)

# Add content to panel 2 (if needed)
content_2 = text_frame_2.add_paragraph()
content_2.text =  "This is the left panel with bullet points:\n\n" \
            "• Bullet point 1\n" \
            "• Bullet point 2\n" \
            "• Bullet point 3\n"



## Add long text with bullet points to the panels
#left_text = "This is the left panel with bullet points:\n\n" \
            #"• Bullet point 1\n" \
            #"• Bullet point 2\n" \
            #"• Bullet point 3\n"
#right_text = "This is the right panel with bullet points:\n\n" \
             #"• Bullet point A\n" \
             #"• Bullet point B\n" \
             #"• Bullet point C\n"

#left_column.text = left_text
#right_column.text = right_text

## Customize the bullet points' styling (font size and type)
#for paragraph in left_column.text_frame.paragraphs:
    #for run in paragraph.runs:
        #run.font.size = Pt(16)
#for paragraph in right_column.text_frame.paragraphs:
    #for run in paragraph.runs:
        #run.font.size = Pt(16)








# ===================================================================== 




# Add Pie Chart
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_3.shapes.title
title.text = "Pie Chart"
# Add your pie chart here (e.g., using Matplotlib or another library)

# Add a Table
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_4.shapes.title
title.text = "Table"
# Add a table using a table object from python-pptx

# Add Two Panels Text
slide_5 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_5.shapes.title
title.text = "Two Panels Text"
left_column = slide_5.shapes.add_textbox(Inches(1), Inches(1), Inches(3.5), Inches(3))
left_text_frame = left_column.text_frame
left_text_frame.text = "Left Panel Text"
left_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

right_column = slide_5.shapes.add_textbox(Inches(4.5), Inches(1), Inches(3.5), Inches(3))
right_text_frame = right_column.text_frame
right_text_frame.text = "Right Panel Text"
right_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add an Interactive Chart
slide_6 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_6.shapes.title
title.text = "Interactive Chart"
# Add an interactive chart using a chart object from python-pptx

# Add Summary
slide_7 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_7.shapes.title
title.text = "Summary"
content = slide_7.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
content_text_frame = content.text_frame
content_text_frame.text = "This prs demonstrates how to create a PowerPoint prs using python-pptx with various elements."

# Add Thank You Slide
slide_8 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_8.shapes.title
title.text = "Thank You!"

# Save the prs
prs.save('python_pptx_prs.pptx')
