 
# pip install python-docx
# pip install -U pypiwin32 OR pip install python-pptx pywin32


image_path =   "/home/godfrey/Documents/Python-PPTX/Image.jpg"


from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
PRES_ = Presentation()

logo_path = "/home/godfrey/Documents/Python-PPTX/Image.jpg"

 
# Slide 0: Cover Page
slide_1 = PRES_.slides.add_slide(PRES_.slide_layouts[5])  # Blank Slide

# Add a full-screen image to the cover page
left_inch = Inches(0)
top_inch = Inches(0)
width_inch = Inches(10)
height_inch = Inches(7.5)
image = slide_1.shapes.add_picture(logo_path, left_inch, top_inch, width_inch, height_inch)

# Add the title text
#                                         xxxx,(0=top, 7=down), (0=left, 7=right),xxxxx
title = slide_1.shapes.add_textbox(Inches(0.5), Inches(7), Inches(2), Inches(1))
title_frame = title.text_frame
title_frame.text = "Test Report" 
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add the name and date at the bottom left
name_date = slide_1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(4), Inches(1))
name_date_frame = name_date.text_frame
from datetime import datetime
name_date_frame.text = "Prepared by: Mosotho MG\nDate:"+str(datetime.now())  
name_date_frame.paragraphs[0].font.size = Pt(18)



# =================================== Slide 1 =================================== #
slide_layout = PRES_.slide_layouts[5]  # Using the blank slide layout
slide = PRES_.slides.add_slide(slide_layout)
# Add the image to the slide and center it
left = (PRES_.slide_width - Inches(3)) / 2  # Center the image horizontally
top = (PRES_.slide_height - Inches(3)) / 2  # Center the image vertically
width = Inches(3)
height = Inches(3)
picture = slide.shapes.add_picture(image_path, left, top, width, height)
# Add a title to the slide
title_text = 'Your Title Here'
title_shape = slide.shapes.title
title_shape.text = title_text
title_shape.text_frame.text = title_text
title_shape.text_frame.paragraphs[0].font.size = Pt(24)  # Adjust the font size as needed
# Add a caption to the slide
caption_text = 'Your Image Caption Here'
caption_shape = slide.shapes.add_textbox(left, top + height + Inches(0.5), width, Inches(1))
caption_frame = caption_shape.text_frame
caption_frame.text = caption_text
caption_frame.paragraphs[0].font.size = Pt(16)  # Adjust the font size as needed

# =================================== Slide 2 =================================== #
slide_layout = PRES_.slide_layouts[5]  # Using the blank slide layout
slide = PRES_.slides.add_slide(slide_layout)
# Add the image to the slide and center it
left = (PRES_.slide_width - Inches(3)) / 2  # Center the image horizontally
top = (PRES_.slide_height - Inches(3)) / 2  # Center the image vertically
width = Inches(3)
height = Inches(3)
picture = slide.shapes.add_picture(image_path, left, top, width, height)
# Add a title to the slide
title_text = 'Your Title Here'
title_shape = slide.shapes.title
title_shape.text = title_text
title_shape.text_frame.text = title_text
title_shape.text_frame.paragraphs[0].font.size = Pt(24)  # Adjust the font size as needed
# Add a caption to the slide
caption_text = 'Your Image Caption Here'
caption_shape = slide.shapes.add_textbox(left, top + height + Inches(0.5), width, Inches(1))
caption_frame = caption_shape.text_frame
caption_frame.text = caption_text
caption_frame.paragraphs[0].font.size = Pt(16)  # Adjust the font size as needed

# =================================== Slide 3 =================================== #
slide_layout = PRES_.slide_layouts[5]  # Using the blank slide layout
slide = PRES_.slides.add_slide(slide_layout)
# Add the image to the slide and center it
left = (PRES_.slide_width - Inches(3)) / 2  # Center the image horizontally
top = (PRES_.slide_height - Inches(3)) / 2  # Center the image vertically
width = Inches(3)
height = Inches(3)
picture = slide.shapes.add_picture(image_path, left, top, width, height)
# Add a title to the slide
title_text = 'Your Title Here'
title_shape = slide.shapes.title
title_shape.text = title_text
title_shape.text_frame.text = title_text
title_shape.text_frame.paragraphs[0].font.size = Pt(24)  # Adjust the font size as needed
# Add a caption to the slide
caption_text = 'Your Image Caption Here'
caption_shape = slide.shapes.add_textbox(left, top + height + Inches(0.5), width, Inches(1))
caption_frame = caption_shape.text_frame
caption_frame.text = caption_text
caption_frame.paragraphs[0].font.size = Pt(16)  # Adjust the font size as needed



excel_file_path = "/home/godfrey/Documents/Python-PPTX/TableX.xlsx"
#from pptx import Presentation
#from pptx.util import Inches
#import io
## Create a presentation object
#from pptx import Presentation
#from pptx.util import Inches 
#title_slide = PRES_.slides.add_slide(PRES_.slide_layouts[5])
#title_slide.shapes.add_ole_object(excel_file_path, 'Excel.Sheet', left = Inches(1), top = Inches(3))
 
#import win32com.client
## initialize Powerpoint and get slide object
#ppt_app = win32com.client.Dispatch("PowerPoint.Application")
#ppt_presentation = ppt_app.Presentations.Add(True)
#ppt_presentation.Slides.Add(1, 12)
#ppt_slide = ppt_presentation.Slides(1)

## load Excel file and get worksheet object
#excel_app = win32com.client.Dispatch("Excel.Application")
#workbook = excel_app.Workbooks.Open(Filename=excel_file_path, ReadOnly=1)
#worksheet = workbook.Sheets(1)
#worksheet.Range("A1:H8").Copy()  # select cells and copy to clipboard
## paste cells to Powerpoint slide
#ppt_slide.Shapes.PasteSpecial(DataType=0, Link=False)






import os
import shutil
import tempfile
from pptx import Presentation
from win32com.client import Dispatch




 
# Save the presentation
PRES_.save('example.pptx')





















































#from pptx import Presentation
#import os
#import shutil
#import tempfile
#from win32com.client import Dispatch

#def embed_excel_slide(prs, excel_file_path):
    ## Create a temporary directory to extract contents
    #temp_dir = tempfile.mkdtemp()

    #try:
        ## Extract contents of Excel file to temporary directory
        #shutil.copy(excel_file_path, temp_dir)
        #excel_file_name = os.path.basename(excel_file_path)
        #temp_excel_path = os.path.join(temp_dir, excel_file_name)

        ## Create a new presentation and add a slide
        #slide_layout = prs.slide_layouts[5]  # Blank slide layout
        #slide = prs.slides.add_slide(slide_layout)

        ## Embed the Excel file into the slide
        #left = top = 0
        #width = height = 1  # Set initial size (will be adjusted later)
        #ole_obj = slide.shapes.add_ole_object(
            #left, top, width, height, temp_excel_path, icon_file='excel.ico'
        #)

        ## Get the embedded object's frame and adjust its size
        #frame = ole_obj.ole_object
        #frame.width = int(prs.slide_width * Inches(0.8))
        #frame.height = int(prs.slide_height * Inches(0.8))

    #finally:
        ## Clean up the temporary directory
        #shutil.rmtree(temp_dir)

## Create a presentation object
#presentation = Presentation()

## Specify the path to your Excel file
#excel_file_path = 'path/to/your/excel_file.xlsx'

## Embed Excel file in the presentation
#embed_excel_slide(presentation, excel_file_path)

## Save the presentation
#presentation.save('example.pptx')












#import matplotlib.pyplot as plt
#from pptx import Presentation
#from pptx.util import Inches
#import tempfile
#import shutil

#def create_python_figure():
    ## Create a simple matplotlib plot
    #fig, ax = plt.subplots()
    #ax.plot([1, 2, 3, 4], [10, 15, 7, 12])
    #ax.set_xlabel('X-axis')
    #ax.set_ylabel('Y-axis')
    #ax.set_title('Matplotlib Plot')

    ## Save the plot as an image file
    #image_path = 'path/to/your/figure.png'
    #plt.savefig(image_path, format='png')

    #return image_path

#def embed_figure_in_presentation(prs, image_path):
    ## Add a slide to the presentation
    #slide_layout = prs.slide_layouts[5]  # Blank slide layout
    #slide = prs.slides.add_slide(slide_layout)

    ## Add the image to the slide
    #left = top = Inches(1)
    #width = height = Inches(4)
    #picture = slide.shapes.add_picture(image_path, left, top, width, height)

## Create a PowerPoint presentation object
#presentation = Presentation()

## Create a Python-generated figure and save it as an image
#figure_path = create_python_figure()

## Embed the figure into the presentation
#embed_figure_in_presentation(presentation, figure_path)

## Save the presentation
#presentation.save('example_with_python_figure.pptx')

## Cleanup: Remove the temporary figure file
#shutil.remove(figure_path)
